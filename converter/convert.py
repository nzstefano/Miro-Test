# converter/convert.py
import argparse
import json
import re
import html
import os
import tempfile
from typing import Any, Dict, Optional, Tuple, List

import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Accurate label width/height using Pillow
try:
    from PIL import ImageFont, Image, ImageDraw
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False


# -----------------------------
# Tunables (easy to tweak)
# -----------------------------
DPI = 96.0

# Label tuning
LABEL_HEIGHT_THRESHOLD_PX = 28   # small boxes w/ bg => treat as label/chip
LABEL_PAD_X_PX = 10
LABEL_PAD_Y_PX = 4
LABEL_HEIGHT_FUDGE_PX = 4        # tiny extra so descenders don’t look clipped
LABEL_MIN_GAP_PX = 6
LABEL_MIN_WIDTH_FRAC = 0.24      # ensures labels “spread horizontally enough”
LABEL_STACK_GAP_PX = 4           # vertical spacing between stacked name/role
LABEL_PAIR_MAX_DY_PX = 64        # max vertical distance to be considered a pair
LABEL_ALIGN_X_TOL_PX = 24        # if no horizontal overlap, allow small x difference

# Image tuning
FIT_IMAGES = True
IMAGE_ENLARGE_MAX = 2.6          # allow a bit more enlargement
IMAGE_SCALE_BIAS = 1.03          # tiny boost
IMAGE_FRAME_MARGIN = 0.02        # 2% margin to slide edges


# -----------------------------
# Helpers
# -----------------------------
def px_to_inches(px: float) -> float:
    return float(px) / DPI

def strip_html_basic(text: str) -> str:
    if not text:
        return ""
    t = re.sub(r"<\s*br\s*/?>", "\n", text, flags=re.I)
    t = re.sub(r"</\s*p\s*>\s*<\s*p\s*>", "\n", t, flags=re.I)
    t = re.sub(r"</?p[^>]*>", "", t, flags=re.I)
    t = re.sub(r"<[^>]+>", "", t)
    return t

def html_has_strong(text: str) -> bool:
    return bool(re.search(r"<\s*(strong|b)\b", text or "", flags=re.I))

def parse_json_str(s: Optional[str]) -> Dict[str, Any]:
    if not s or not isinstance(s, str):
        return {}
    try:
        return json.loads(s)
    except Exception:
        return {}

def int_color_to_rgb(value: Any) -> Optional[Tuple[int, int, int]]:
    try:
        n = int(value)
    except (TypeError, ValueError):
        return None
    if n < 0:
        return None
    return ((n >> 16) & 0xFF, (n >> 8) & 0xFF, n & 0xFF)

def rgb_tuple_to_hex(rgb: Tuple[int, int, int]) -> str:
    return "#{:02X}{:02X}{:02X}".format(*rgb)

def hex_to_rgb(hex_color: Optional[str]) -> Optional[Tuple[int, int, int]]:
    if not hex_color:
        return None
    h = hex_color.strip().lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return None
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def get_scale(jdata: Dict[str, Any]) -> float:
    s = float(((jdata.get("scale") or {}).get("scale")) or 1.0)
    rs = float(jdata.get("relativeScale") or 1.0)
    return s * rs

def center_to_topleft(cx: float, cy: float, w: float, h: float) -> Tuple[float, float]:
    return float(cx) - float(w) / 2.0, float(cy) - float(h) / 2.0

def map_align(ta: Optional[str]) -> Optional[PP_ALIGN]:
    if not ta:
        return None
    ta = ta.lower()
    return {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}.get(ta)

def measure_text_width_px(text: str, font_name: Optional[str], font_size_pt: float) -> float:
    """Accurate label width in px using Pillow; fallback to a generous estimate."""
    if not text:
        return 0.0
    if not PIL_AVAILABLE:
        return len(text) * font_size_pt * (96.0/72.0) * 0.65  # slightly generous

    px_size = max(1, int(round(font_size_pt * (96.0/72.0))))
    font = None
    candidates = []
    if font_name:
        candidates.append(font_name)
    candidates += ["Roobert.ttf", "NotoSans-Regular.ttf", "DejaVuSans.ttf", "Arial.ttf", "LiberationSans-Regular.ttf"]
    for name in candidates:
        try:
            font = ImageFont.truetype(name, px_size)
            break
        except Exception:
            continue
    if font is None:
        font = ImageFont.load_default()

    img = Image.new("RGB", (1, 1))
    draw = ImageDraw.Draw(img)
    bbox = draw.textbbox((0, 0), text, font=font)
    width = (bbox[2] - bbox[0]) if bbox else 0
    return width * 1.08  # safety margin

def measure_text_bbox_px(text: str, font_name: Optional[str], font_size_pt: float) -> Tuple[float, float]:
    """Return (width, height) in px using Pillow; fallback to safe estimates."""
    if not text:
        return 0.0, 0.0
    if not PIL_AVAILABLE:
        width = len(text) * font_size_pt * (96.0/72.0) * 0.65
        height = font_size_pt * (96.0/72.0) * 1.20
        return width, height

    px_size = max(1, int(round(font_size_pt * (96.0/72.0))))
    font = None
    candidates = []
    if font_name:
        candidates.append(font_name)
    candidates += ["Roobert.ttf", "NotoSans-Regular.ttf", "DejaVuSans.ttf", "Arial.ttf", "LiberationSans-Regular.ttf"]
    for name in candidates:
        try:
            font = ImageFont.truetype(name, px_size)
            break
        except Exception:
            continue
    if font is None:
        font = ImageFont.load_default()

    img = Image.new("RGB", (1, 1))
    draw = ImageDraw.Draw(img)
    bbox = draw.textbbox((0, 0), text, font=font)  # (x0, y0, x1, y1)
    if not bbox:
        return 0.0, px_size
    w = (bbox[2] - bbox[0]) * 1.08  # small safety margins
    h = (bbox[3] - bbox[1]) * 1.05
    return w, h

def horiz_intersect(a: Dict[str, float], b: Dict[str, float]) -> bool:
    return not (a["left"] + a["width"] <= b["left"] or b["left"] + b["width"] <= a["left"])


# -----------------------------
# Parse Miro export -> model
# -----------------------------
def parse_miro_export(data: Dict[str, Any]) -> Dict[str, Any]:
    widgets: List[Dict[str, Any]] = (data.get("content") or {}).get("widgets") or []

    # Frames (ordered)
    frames: List[Dict[str, Any]] = []
    for w in widgets:
        c = w.get("canvasedObjectData") or {}
        if c.get("type") == "frame":
            j = parse_json_str(c.get("json"))
            size = j.get("size") or {}
            style = parse_json_str(j.get("style"))
            frames.append({
                "id": w.get("id"),
                "order": j.get("presentationOrder", ""),
                "size": {"width": float(size.get("width", 1280.0)),
                         "height": float(size.get("height", 720.0))},
                "style": style,
            })
    frames.sort(key=lambda f: f["order"])
    slide_size_px = frames[0]["size"].copy() if frames else {"width": 1280.0, "height": 720.0}

    slides: List[Dict[str, Any]] = []
    frame_iter = frames if frames else [{"id": None, "style": {}, "size": slide_size_px}]

    for frame in frame_iter:
        frame_id = frame["id"]
        frame_style = frame.get("style") or {}
        bg_rgb = int_color_to_rgb(frame_style.get("bc"))
        bg_hex = rgb_tuple_to_hex(bg_rgb) if bg_rgb else None

        items: List[Dict[str, Any]] = []

        for w in widgets:
            c = w.get("canvasedObjectData") or {}
            wtype = c.get("type")
            jdata = parse_json_str(c.get("json"))

            if frame_id is not None:
                parent_id = ((jdata.get("_parent") or {}).get("id"))
                if parent_id != frame_id:
                    continue

            # geometry (Miro uses center coords)
            pos = ((jdata.get("_position") or {}).get("offsetPx")) or {}
            size = (jdata.get("size") or {})
            cx = float(pos.get("x", 0.0))
            cy = float(pos.get("y", 0.0))
            base_w = float(size.get("width", 200.0))
            base_h = float(size.get("height", 100.0))
            scale = get_scale(jdata)

            wpx = base_w * scale
            hpx = base_h * scale
            left, top = center_to_topleft(cx, cy, wpx, hpx)

            if wtype == "text":
                raw = jdata.get("text", "")
                clean = strip_html_basic(raw)
                style = parse_json_str(jdata.get("style"))
                fs = style.get("fs")
                # NOTE: we NO LONGER trust style["b"] because it can be noisy;
                # bold only if HTML contains <strong>/<b>
                bold_html = html_has_strong(raw)
                i = bool(style.get("i")) if "i" in style else False
                tc = int_color_to_rgb(style.get("tc"))
                ta = style.get("ta")
                ffn = style.get("ffn")

                # detect label (bg color present + small height)
                is_label = (style.get("bc") is not None and style.get("bc") != -1 and base_h <= LABEL_HEIGHT_THRESHOLD_PX)

                item: Dict[str, Any] = {
                    "type": "text",
                    "text": html.unescape(clean),
                    "left": left, "top": top, "width": wpx, "height": hpx,
                    "isLabel": is_label,
                }
                if fs:
                    item["fontSize"] = int(round(float(fs) * scale))
                if bold_html:
                    item["bold"] = True
                if i:
                    item["italic"] = True
                if tc:
                    item["color"] = rgb_tuple_to_hex(tc)
                if ta:
                    item["align"] = ta
                if ffn:
                    item["font"] = ffn

                if is_label:
                    fill_rgb = int_color_to_rgb(style.get("bc"))
                    line_rgb = int_color_to_rgb(style.get("sc"))
                    brw = style.get("brw")
                    if fill_rgb: item["boxFill"] = rgb_tuple_to_hex(fill_rgb)
                    if line_rgb: item["boxLine"] = rgb_tuple_to_hex(line_rgb)
                    if isinstance(brw, (int, float)): item["boxLineWidth"] = float(brw) * scale
                    item["noWrap"] = True  # single-line chip

                    # Use actual text metrics for BOTH width & height
                    fs_pt = item.get("fontSize") or 14
                    mw, mh = measure_text_bbox_px(item["text"], item.get("font"), fs_pt)

                    # Width: text + horizontal padding (keep center)
                    needed_w = mw + 2*LABEL_PAD_X_PX
                    if needed_w > item["width"]:
                        cx_lbl = item["left"] + item["width"]/2.0
                        item["width"] = needed_w
                        item["left"] = cx_lbl - needed_w/2.0

                    # Height: text height + vertical padding + tiny fudge (keep middle)
                    min_h = mh + 2*LABEL_PAD_Y_PX + LABEL_HEIGHT_FUDGE_PX
                    if item["height"] < min_h:
                        delta = (min_h - item["height"]) / 2.0
                        item["top"] -= delta
                        item["height"] = min_h

                items.append(item)

            elif wtype == "image":
                img = jdata.get("image") or {}
                link = img.get("externalLink")
                if not link:
                    continue

                resource = jdata.get("resource") or {}
                crop = jdata.get("crop") or {}
                size_w = (jdata.get("size") or {}).get("width")
                size_h = (jdata.get("size") or {}).get("height")

                # prefer explicit size, then crop, then resource — all × scale
                if size_w and size_h:
                    wpx2 = float(size_w) * scale
                    hpx2 = float(size_h) * scale
                else:
                    cw = float(crop.get("width") or 0.0)
                    ch = float(crop.get("height") or 0.0)
                    if cw > 0 and ch > 0:
                        wpx2 = cw * scale
                        hpx2 = ch * scale
                    else:
                        rw = float(resource.get("width") or 0.0)
                        rh = float(resource.get("height") or 0.0)
                        if rw > 0 and rh > 0:
                            wpx2 = rw * scale
                            hpx2 = rh * scale
                        else:
                            wpx2, hpx2 = wpx, hpx

                left2, top2 = center_to_topleft(cx, cy, wpx2, hpx2)
                items.append({
                    "type": "image",
                    "url": link,
                    "left": left2, "top": top2, "width": wpx2, "height": hpx2,
                })

            elif wtype in {"shape", "stickynote", "sticker", "container", "card"}:
                style = parse_json_str(jdata.get("style"))
                fill_rgb = int_color_to_rgb(style.get("bc"))
                line_rgb = int_color_to_rgb(style.get("sc"))
                line_width = style.get("brw")
                corner = style.get("brr", 0.0)

                item = {
                    "type": "shape",
                    "shape": "round_rect" if (isinstance(corner, (int, float)) and corner > 0) else "rectangle",
                    "left": left, "top": top, "width": wpx, "height": hpx,
                }
                if fill_rgb: item["fill"] = rgb_tuple_to_hex(fill_rgb)
                if line_rgb: item["line"] = rgb_tuple_to_hex(line_rgb)
                if isinstance(line_width, (int, float)): item["lineWidth"] = float(line_width) * scale
                items.append(item)

        # Post-process per slide: label tidy/alignment & image enlarge
        items = tidy_labels_and_images(items, frame["size"]["width"], frame["size"]["height"])
        slides.append({"bgColor": bg_hex, "items": items})

    return {"slideSizePx": slide_size_px, "slides": slides}


def rects_overlap(a: Dict[str, float], b: Dict[str, float]) -> bool:
    return not (
        a["left"] + a["width"] <= b["left"] or
        b["left"] + b["width"] <= a["left"] or
        a["top"] + a["height"] <= b["top"] or
        b["top"] + b["height"] <= a["top"]
    )

def tidy_labels_and_images(items: List[Dict[str, Any]], slide_w: float, slide_h: float) -> List[Dict[str, Any]]:
    # --- LABELS ---
    labels = [i for i in items if i.get("type") == "text" and i.get("isLabel")]
    labels.sort(key=lambda t: (t["top"], t["left"]))

    # 1) Pair alignment: stack pairs (name/role) with same left & width
    used = set()
    for i, a in enumerate(labels):
        if id(a) in used:
            continue
        # find a partner just below a (closest dy)
        best = None
        best_dy = 1e9
        for j in range(i + 1, len(labels)):
            b = labels[j]
            dy = b["top"] - a["top"]
            if dy < 0 or dy > LABEL_PAIR_MAX_DY_PX:
                continue
            # accept if horizontally intersect OR near-left within tolerance
            if not horiz_intersect(a, b):
                if abs(a["left"] - b["left"]) > max(LABEL_ALIGN_X_TOL_PX, 0.05 * slide_w):
                    continue
            if dy < best_dy:
                best = b
                best_dy = dy
        if best:
            common_left = min(a["left"], best["left"])
            common_width = max(a["width"], best["width"], slide_w * LABEL_MIN_WIDTH_FRAC)
            a["left"] = best["left"] = common_left
            a["width"] = best["width"] = common_width

            # order by top, then stack with a tidy gap
            top_label = a if a["top"] <= best["top"] else best
            bottom_label = best if top_label is a else a
            bottom_label["top"] = top_label["top"] + top_label["height"] + LABEL_STACK_GAP_PX

            used.add(id(a)); used.add(id(best))

    # 2) Enforce min width & de-overlap vertically
    placed: List[Dict[str, Any]] = []
    for lb in labels:
        min_w = max(lb["width"], slide_w * LABEL_MIN_WIDTH_FRAC)
        if lb["width"] < min_w:
            cx = lb["left"] + lb["width"]/2.0
            lb["width"] = min_w
            lb["left"] = cx - min_w/2.0

        for prev in placed:
            if not (lb["left"] + lb["width"] <= prev["left"] or prev["left"] + prev["width"] <= lb["left"]):
                while rects_overlap(lb, prev):
                    lb["top"] = prev["top"] + prev["height"] + LABEL_MIN_GAP_PX
        placed.append(lb)

    # --- IMAGES ---
    if FIT_IMAGES:
        # Reserve header space only where labels actually are (near top)
        reserved_top = 0.0
        for lb in labels:
            if lb["top"] < (slide_h * 0.35):
                reserved_top = max(reserved_top, lb["top"] + lb["height"] + LABEL_MIN_GAP_PX)

        margin_w = slide_w * IMAGE_FRAME_MARGIN
        avail_w = max(1.0, slide_w - 2 * margin_w)
        avail_h = max(1.0, (slide_h - reserved_top) - (slide_h * IMAGE_FRAME_MARGIN))

        for it in items:
            if it.get("type") != "image":
                continue

            # Fit by free area
            s_fit_w = avail_w / it["width"] if it["width"] > 0 else 1.0
            s_fit_h = avail_h / it["height"] if it["height"] > 0 else 1.0
            s_fit = min(s_fit_w, s_fit_h)

            # Keep within edges given the current center
            cx = it["left"] + it["width"] / 2.0
            cy = it["top"] + it["height"] / 2.0
            max_half_w = min(cx - margin_w, (slide_w - margin_w) - cx)
            max_half_h = min((cy - reserved_top), (slide_h - (slide_h * IMAGE_FRAME_MARGIN)) - cy)
            max_half_w = max(max_half_w, 1.0)
            max_half_h = max(max_half_h, 1.0)

            s_center_w = max_half_w / (it["width"]/2.0)
            s_center_h = max_half_h / (it["height"]/2.0)
            s_center = min(s_center_w, s_center_h)

            s = min(max(1.0, s_fit), s_center, IMAGE_ENLARGE_MAX) * IMAGE_SCALE_BIAS
            if s > 1.0:
                new_w = it["width"] * s
                new_h = it["height"] * s
                it["left"], it["top"] = center_to_topleft(cx, cy, new_w, new_h)
                it["width"], it["height"] = new_w, new_h

    return items


# -----------------------------
# Render to PPTX
# -----------------------------
def add_frame_background(slide, slide_meta: Dict[str, Any], slide_width_px: float, slide_height_px: float):
    bg = slide_meta.get("bgColor")
    if not bg:
        return
    rgb = hex_to_rgb(bg)
    if not rgb:
        return
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(px_to_inches(slide_width_px)),
        Inches(px_to_inches(slide_height_px)),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*rgb)
    rect.line.fill.background()

def add_textbox(slide, item: Dict[str, Any]):
    x = Inches(px_to_inches(item["left"]))
    y = Inches(px_to_inches(item["top"]))
    w = Inches(px_to_inches(item["width"]))
    h = Inches(px_to_inches(item["height"]))

    box = slide.shapes.add_textbox(x, y, w, h)

    # Label background/border + paddings
    box_fill = item.get("boxFill")
    if box_fill:
        rgb = hex_to_rgb(box_fill)
        if rgb:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*rgb)
    box_line = item.get("boxLine")
    if box_line:
        rgb = hex_to_rgb(box_line)
        if rgb:
            box.line.color.rgb = RGBColor(*rgb)
    blw = item.get("boxLineWidth")
    if isinstance(blw, (int, float)):
        box.line.width = Pt(blw)

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = not item.get("noWrap", False)
    tf.vertical_anchor = MSO_ANCHOR.TOP if item.get("isLabel") else MSO_ANCHOR.MIDDLE

    if item.get("isLabel"):
        tf.margin_left   = Inches(px_to_inches(LABEL_PAD_X_PX))
        tf.margin_right  = Inches(px_to_inches(LABEL_PAD_X_PX))
        tf.margin_top    = Inches(px_to_inches(LABEL_PAD_Y_PX))
        tf.margin_bottom = Inches(px_to_inches(LABEL_PAD_Y_PX))

    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    try:
        p.line_spacing = 1.0
    except Exception:
        pass

    run = p.add_run()
    run.text = item.get("text", "")
    font = run.font

    fs = item.get("fontSize")
    if isinstance(fs, (int, float)):
        font.size = Pt(fs)
    if item.get("bold"):
        font.bold = True
    if item.get("italic"):
        font.italic = True
    if item.get("font"):
        font.name = str(item["font"])

    color = item.get("color")
    rgb = hex_to_rgb(color) if color else None
    if rgb:
        font.color.rgb = RGBColor(*rgb)

    align = map_align(item.get("align"))
    if align is not None:
        p.alignment = align

def add_shape(slide, item: Dict[str, Any]):
    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    }
    st = shape_map.get(str(item.get("shape", "rectangle")).lower(), MSO_SHAPE.RECTANGLE)
    x = Inches(px_to_inches(item["left"]))
    y = Inches(px_to_inches(item["top"]))
    w = Inches(px_to_inches(item["width"]))
    h = Inches(px_to_inches(item["height"]))
    shp = slide.shapes.add_shape(st, x, y, w, h)

    fill_color = hex_to_rgb(item.get("fill"))
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*fill_color)

    line_color = hex_to_rgb(item.get("line"))
    if line_color:
        shp.line.color.rgb = RGBColor(*line_color)

    lw = item.get("lineWidth")
    if isinstance(lw, (int, float)):
        shp.line.width = Pt(lw)

def add_image(slide, item: Dict[str, Any]):
    try:
        resp = requests.get(item["url"], stream=True, timeout=15)
        resp.raise_for_status()
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".img")
        with open(tmp.name, "wb") as f:
            for chunk in resp.iter_content(1024):
                f.write(chunk)
        slide.shapes.add_picture(
            tmp.name,
            Inches(px_to_inches(item["left"])),
            Inches(px_to_inches(item["top"])),
            Inches(px_to_inches(item["width"])),
            Inches(px_to_inches(item["height"])),
        )
    except Exception as e:
        print(f"[warn] image load failed: {item.get('url')} -> {e}")

def build_presentation(model: Dict[str, Any]) -> Presentation:
    prs = Presentation()
    slide_w_px = float((model.get("slideSizePx") or {}).get("width") or 1280.0)
    slide_h_px = float((model.get("slideSizePx") or {}).get("height") or 720.0)
    prs.slide_width = Inches(px_to_inches(slide_w_px))
    prs.slide_height = Inches(px_to_inches(slide_h_px))

    for s in model.get("slides", []):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        add_frame_background(slide, s, slide_w_px, slide_h_px)
        for it in s.get("items", []):
            t = it.get("type")
            if t == "text":
                add_textbox(slide, it)
            elif t == "shape":
                add_shape(slide, it)
            elif t == "image":
                add_image(slide, it)
    return prs


# -----------------------------
# Entry
# -----------------------------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True, help="Path to input JSON")
    ap.add_argument("--output", required=True, help="Path to output PPTX")
    args = ap.parse_args()

    out_dir = os.path.dirname(args.output) or "."
    os.makedirs(out_dir, exist_ok=True)

    with open(args.input, "r") as f:
        data = json.load(f)

    model = parse_miro_export(data) if ("content" in data and "widgets" in (data["content"] or {})) else data
    prs = build_presentation(model)
    prs.save(args.output)
    print(f"Saved {args.output}")

if __name__ == "__main__":
    main()
